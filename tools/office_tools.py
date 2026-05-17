"""
Office file tools: Word (.docx), Excel (.xlsx), PowerPoint (.pptx) の読み書き。

依存: python-docx, openpyxl, python-pptx
インストール: pip install python-docx openpyxl python-pptx
"""

from pathlib import Path
from tools.file_tools import _resolve_safe_path


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------

def read_docx(path: str) -> dict:
    """
    Word ファイル (.docx) を読み込み、テキストを返します。

    path: workspace 相対パス (例: docs/report.docx)
    """
    try:
        from docx import Document
    except ImportError:
        return {"error": "python-docx がインストールされていません。run_command('pip install python-docx') でインストールしてください。"}

    try:
        target = _resolve_safe_path(path)
        if not target.exists():
            return {"error": f"ファイルが見つかりません: {path}"}

        doc = Document(str(target))
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            level = ""
            if style.startswith("Heading"):
                try:
                    n = int(style.split()[-1])
                    level = "#" * n + " "
                except ValueError:
                    level = "# "
            paragraphs.append({"level": level, "text": text, "style": style})

        full_text = "\n".join((p["level"] + p["text"]) for p in paragraphs)
        return {
            "text": full_text,
            "paragraphs": paragraphs,
            "paragraph_count": len(paragraphs),
            "section_count": len(doc.sections),
            "path": str(target),
        }
    except Exception as e:
        return {"error": f"Word ファイル読み込みエラー: {e}"}


def write_docx(path: str, content: str, title: str = "") -> dict:
    """
    Word ファイル (.docx) を作成・上書きします。
    content は Markdown 風テキスト（# 見出し、## 小見出し、通常段落）。

    path: workspace 相対パス (例: docs/report.docx)
    content: Markdown 風テキスト
    title: ドキュメントタイトル（省略可）
    """
    try:
        from docx import Document
    except ImportError:
        return {"error": "python-docx がインストールされていません。run_command('pip install python-docx') でインストールしてください。"}

    try:
        target = _resolve_safe_path(path)
        target.parent.mkdir(parents=True, exist_ok=True)

        doc = Document()
        if title:
            doc.add_heading(title, level=0)

        count = 0
        for line in content.splitlines():
            stripped = line.rstrip()
            if stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            elif stripped == "":
                continue
            else:
                doc.add_paragraph(stripped)
            count += 1

        doc.save(str(target))
        return {"path": str(target), "paragraphs_written": count, "error": None}
    except Exception as e:
        return {"error": f"Word ファイル書き込みエラー: {e}"}


def edit_docx(path: str, old_text: str, new_text: str) -> dict:
    """
    Word ファイル内の指定テキストを置換します（段落単位）。

    path: workspace 相対パス
    old_text: 置換前のテキスト
    new_text: 置換後のテキスト
    """
    try:
        from docx import Document
    except ImportError:
        return {"error": "python-docx がインストールされていません。"}

    try:
        target = _resolve_safe_path(path)
        if not target.exists():
            return {"error": f"ファイルが見つかりません: {path}"}

        doc = Document(str(target))
        count = 0
        for para in doc.paragraphs:
            if old_text in para.text:
                for run in para.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                        count += 1

        if count == 0:
            return {"error": f"'{old_text}' が見つかりませんでした。read_docx でテキストを確認してください。"}

        doc.save(str(target))
        return {"replaced_count": count, "path": str(target), "error": None}
    except Exception as e:
        return {"error": f"Word ファイル編集エラー: {e}"}


# ---------------------------------------------------------------------------
# Excel (.xlsx)
# ---------------------------------------------------------------------------

def read_xlsx(path: str, sheet: str = None, max_rows: int = 200) -> dict:
    """
    Excel ファイル (.xlsx) を読み込み、シートのデータを返します。

    path: workspace 相対パス (例: data/sales.xlsx)
    sheet: シート名（省略時は最初のシート）
    max_rows: 最大読み込み行数（デフォルト 200）
    """
    try:
        import openpyxl
    except ImportError:
        return {"error": "openpyxl がインストールされていません。run_command('pip install openpyxl') でインストールしてください。"}

    try:
        target = _resolve_safe_path(path)
        if not target.exists():
            return {"error": f"ファイルが見つかりません: {path}"}

        wb = openpyxl.load_workbook(str(target), read_only=True, data_only=True)
        sheet_names = wb.sheetnames
        ws = wb[sheet] if sheet and sheet in sheet_names else wb.active
        actual_sheet = ws.title

        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows:
                break
            rows.append([str(cell) if cell is not None else "" for cell in row])
        wb.close()

        headers = rows[0] if rows else []
        data_rows = rows[1:] if len(rows) > 1 else []

        return {
            "sheet_name": actual_sheet,
            "sheets": sheet_names,
            "headers": headers,
            "rows": data_rows,
            "row_count": len(data_rows),
            "error": None,
        }
    except Exception as e:
        return {"error": f"Excel 読み込みエラー: {e}"}


def write_xlsx(path: str, data: list, sheet: str = "Sheet1", headers: list = None) -> dict:
    """
    Excel ファイル (.xlsx) を作成・上書きします。

    path: workspace 相対パス (例: output/result.xlsx)
    data: 行データのリスト（例: [["Alice", 30], ["Bob", 25]]）
    sheet: シート名（デフォルト: Sheet1）
    headers: ヘッダー行（省略可）
    """
    try:
        import openpyxl
        from openpyxl.styles import Font
    except ImportError:
        return {"error": "openpyxl がインストールされていません。run_command('pip install openpyxl') でインストールしてください。"}

    try:
        target = _resolve_safe_path(path)
        target.parent.mkdir(parents=True, exist_ok=True)

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet

        row_num = 1
        if headers:
            ws.append(headers)
            for cell in ws[row_num]:
                cell.font = Font(bold=True)
            row_num += 1

        for row in data:
            ws.append(row)

        wb.save(str(target))
        return {"path": str(target), "rows_written": len(data), "error": None}
    except Exception as e:
        return {"error": f"Excel 書き込みエラー: {e}"}


def edit_xlsx(path: str, sheet: str = None, row: int = None, col: int = None,
              cell: str = None, value: str = None) -> dict:
    """
    Excel ファイルの特定セルを編集します。

    path: workspace 相対パス
    sheet: シート名（省略時は最初のシート）
    row: 行番号（1始まり）。cell 指定時は不要
    col: 列番号（1始まり）。cell 指定時は不要
    cell: セルアドレス（例: "B3"）
    value: 設定する値
    """
    try:
        import openpyxl
    except ImportError:
        return {"error": "openpyxl がインストールされていません。"}

    try:
        target = _resolve_safe_path(path)
        if not target.exists():
            return {"error": f"ファイルが見つかりません: {path}"}

        wb = openpyxl.load_workbook(str(target))
        sheet_names = wb.sheetnames
        ws = wb[sheet] if sheet and sheet in sheet_names else wb.active

        if cell:
            ws[cell] = value
            cell_ref = cell
        elif row and col:
            ws.cell(row=row, column=col, value=value)
            cell_ref = f"R{row}C{col}"
        else:
            return {"error": "cell または row+col を指定してください"}

        wb.save(str(target))
        return {"path": str(target), "cell": cell_ref, "value": value, "error": None}
    except Exception as e:
        return {"error": f"Excel 編集エラー: {e}"}


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------

def read_pptx(path: str) -> dict:
    """
    PowerPoint ファイル (.pptx) を読み込み、スライドのテキストと画像情報を返します。

    path: workspace 相対パス (例: slides/presentation.pptx)
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx がインストールされていません。run_command('pip install python-pptx') でインストールしてください。"}

    try:
        target = _resolve_safe_path(path)
        if not target.exists():
            return {"error": f"ファイルが見つかりません: {path}"}

        prs = Presentation(str(target))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            image_count = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                # shape_type 13 = PICTURE
                if getattr(shape, "shape_type", None) == 13:
                    image_count += 1
            slides.append({
                "slide_number": i + 1,
                "texts": texts,
                "text": "\n".join(texts),
                "image_count": image_count,
            })

        return {
            "slides": slides,
            "slide_count": len(slides),
            "path": str(target),
            "error": None,
        }
    except Exception as e:
        return {"error": f"PowerPoint 読み込みエラー: {e}"}


def write_pptx(path: str, slides: list, title: str = "") -> dict:
    """
    PowerPoint ファイル (.pptx) を作成・上書きします。

    path: workspace 相対パス (例: output/presentation.pptx)
    slides: スライド定義のリスト。各要素は以下のフィールドを持つ:
      - title: スライドタイトル（省略可）
      - content: 本文テキスト（改行区切り、省略可）
      - image_path: 埋め込む画像の workspace 相対パス（省略可）
        例: "GRAAA/AI_Output_Images/generated_xxx.png"
      ※ content のみ → テキストスライド
      ※ image_path のみ → 画像のみスライド（中央配置）
      ※ content + image_path → 左テキスト・右画像のレイアウト
    title: プレゼンテーション全体のタイトル（最初のスライドに使用、省略可）
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
    except ImportError:
        return {"error": "python-pptx がインストールされていません。run_command('pip install python-pptx') でインストールしてください。"}

    try:
        target = _resolve_safe_path(path)
        target.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        title_layout    = prs.slide_layouts[0]
        content_layout  = prs.slide_layouts[1]
        blank_layout    = prs.slide_layouts[6]

        sw = prs.slide_width   # スライド幅 (EMU)
        sh = prs.slide_height  # スライド高さ (EMU)
        margin = Inches(0.4)

        def _add_image(slide, img_path: str, left, top, width, height):
            img_target = _resolve_safe_path(img_path)
            if not img_target.exists():
                raise FileNotFoundError(f"画像ファイルが見つかりません: {img_path}")
            slide.shapes.add_picture(str(img_target), left, top, width, height)

        count = 0
        if title:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = ""
            count += 1

        for s in slides:
            slide_title   = s.get("title", "")
            # content / text どちらのキーも受け付ける
            slide_content = s.get("content", "") or s.get("text", "")
            # image_path / image どちらのキーも受け付ける
            image_path    = s.get("image_path", "") or s.get("image", "")
            # elements 形式のフォールバック: [{"type":"image","path":"..."}]
            if not image_path:
                for el in s.get("elements", []):
                    if el.get("type") == "image" and el.get("path"):
                        image_path = el["path"]
                        break

            if slide_content and image_path:
                # 左テキスト・右画像レイアウト
                slide = prs.slides.add_slide(blank_layout)
                title_top = margin
                title_h   = Inches(0.7)
                if slide_title:
                    txb = slide.shapes.add_textbox(margin, title_top, sw - margin * 2, title_h)
                    txb.text_frame.text = slide_title
                    txb.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
                    txb.text_frame.paragraphs[0].runs[0].font.bold = True
                content_top = title_top + title_h + Inches(0.1)
                content_h   = sh - content_top - margin
                half_w      = (sw - margin * 3) // 2
                # テキストボックス（左半分）
                txb2 = slide.shapes.add_textbox(margin, content_top, half_w, content_h)
                tf = txb2.text_frame
                tf.word_wrap = True
                for i, line in enumerate(slide_content.splitlines()):
                    if i == 0:
                        tf.paragraphs[0].text = line
                    else:
                        tf.add_paragraph().text = line
                # 画像（右半分）
                img_left = margin * 2 + half_w
                _add_image(slide, image_path, img_left, content_top, half_w, content_h)

            elif image_path:
                # 画像のみスライド（中央配置）
                slide = prs.slides.add_slide(blank_layout)
                if slide_title:
                    title_h = Inches(0.7)
                    txb = slide.shapes.add_textbox(margin, margin, sw - margin * 2, title_h)
                    txb.text_frame.text = slide_title
                    txb.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
                    txb.text_frame.paragraphs[0].runs[0].font.bold = True
                    img_top = margin + title_h + Inches(0.1)
                else:
                    img_top = margin
                img_h = sh - img_top - margin
                img_w = sw - margin * 2
                _add_image(slide, image_path, margin, img_top, img_w, img_h)

            elif slide_content:
                # テキストのみスライド
                slide = prs.slides.add_slide(content_layout)
                slide.shapes.title.text = slide_title
                tf = slide.placeholders[1].text_frame
                tf.text = ""
                for i, line in enumerate(slide_content.splitlines()):
                    if i == 0:
                        tf.paragraphs[0].text = line
                    else:
                        tf.add_paragraph().text = line
            else:
                # タイトルのみスライド
                slide = prs.slides.add_slide(title_layout)
                slide.shapes.title.text = slide_title
            count += 1

        prs.save(str(target))
        return {"path": str(target), "slides_written": count, "error": None}
    except FileNotFoundError as e:
        return {"error": str(e)}
    except Exception as e:
        return {"error": f"PowerPoint 書き込みエラー: {e}"}


def edit_pptx(path: str, slide_number: int, old_text: str, new_text: str) -> dict:
    """
    PowerPoint の特定スライドのテキストを置換します。

    path: workspace 相対パス
    slide_number: スライド番号（1始まり）
    old_text: 置換前のテキスト
    new_text: 置換後のテキスト
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx がインストールされていません。"}

    try:
        target = _resolve_safe_path(path)
        if not target.exists():
            return {"error": f"ファイルが見つかりません: {path}"}

        prs = Presentation(str(target))
        if slide_number < 1 or slide_number > len(prs.slides):
            return {"error": f"スライド番号 {slide_number} は範囲外です（全{len(prs.slides)}枚）"}

        slide = prs.slides[slide_number - 1]
        count = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            count += 1

        if count == 0:
            return {"error": f"スライド{slide_number}に '{old_text}' が見つかりませんでした。"}

        prs.save(str(target))
        return {"replaced_count": count, "path": str(target), "error": None}
    except Exception as e:
        return {"error": f"PowerPoint 編集エラー: {e}"}
